from datetime import datetime
from pathlib import Path
import shutil
from typing import List, Optional
from uuid import uuid4

from fastapi import Depends, FastAPI, File, Form, HTTPException, UploadFile, status
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pptx import Presentation
from dotenv import load_dotenv

from .bedrock import BedrockAgentError, analyze_slide_text
from .config import STORAGE_DIR, UPLOAD_DIR
from .db import get_db, init_db
from .schemas import (
    AnalysisDocument,
    Sermon,
    SlideAnalysis,
    SlideContent,
    SlideDecision,
    SlideDecisionPayload,
    DecisionsDocument,
    Suggestion,
)
from .state import init_sermon_state, load_analysis, load_decisions, save_analysis, save_decisions

app = FastAPI(title="Apologia API", version="0.1.0")
load_dotenv(Path(__file__).resolve().parent.parent / ".env")
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=False,
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.on_event("startup")
def startup_event() -> None:
    init_db()


def _ensure_pptx(file: UploadFile) -> None:
    filename = file.filename or ""
    if not filename.lower().endswith(".pptx"):
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail="Only .pptx files are supported in MVP0.",
        )


def _row_to_sermon(row) -> Sermon:
    created_at = datetime.fromisoformat(row["created_at"])
    relative_path = f"uploads/{row['id']}/{row['original_filename']}"
    return Sermon(
        id=row["id"],
        sermonName=row["sermon_name"],
        seriesName=row["series_name"],
        weekOrDate=row["week_or_date"],
        pastorName=row["pastor_name"],
        status=row["status"],
        filePath=relative_path,
        originalFilename=row["original_filename"],
        createdAt=created_at,
    )


def _extract_slide_text(slide) -> str:
    text_chunks = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = (shape.text or "").strip()
        if text:
            text_chunks.append(text)

    try:
        notes_text = slide.notes_slide.notes_text_frame.text
    except Exception:
        notes_text = ""

    notes_text = (notes_text or "").strip()
    if notes_text:
        text_chunks.append(f"Notes:\n{notes_text}")

    return "\n".join(text_chunks).strip()


def _resolve_upload_path(
    sermon_id: str, file_path: str, original_filename: str
) -> Path:
    stored_path = Path(file_path)
    if stored_path.is_absolute() and stored_path.exists():
        return stored_path

    candidate = UPLOAD_DIR / file_path
    if candidate.exists():
        return candidate

    return UPLOAD_DIR / sermon_id / original_filename


def _ensure_sermon_exists(db, sermon_id: str) -> None:
    row = db.execute(
        """
        SELECT id
        FROM sermons
        WHERE id = ?
        """,
        (sermon_id,),
    ).fetchone()
    if not row:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Not found")


def _get_presentation(db, sermon_id: str) -> Presentation:
    row = db.execute(
        """
        SELECT id, file_path, original_filename
        FROM sermons
        WHERE id = ?
        """,
        (sermon_id,),
    ).fetchone()
    if not row:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Not found")

    file_path = _resolve_upload_path(
        row["id"], row["file_path"], row["original_filename"]
    )
    if not file_path.exists():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="File missing")

    return Presentation(file_path)


def _analyze_text_stub(text: str) -> List[Suggestion]:
    return []


@app.post("/sermons", response_model=Sermon, status_code=status.HTTP_201_CREATED)
async def upload_sermon(
    file: UploadFile = File(...),
    weekOrDate: Optional[str] = Form(None),
    seriesName: Optional[str] = Form(None),
    sermonName: str = Form(...),
    pastorName: Optional[str] = Form(None),
    db=Depends(get_db),
) -> Sermon:
    """
    Store a sermon PPTX file with optional metadata.
    """
    _ensure_pptx(file)

    sermon_id = str(uuid4())
    created_at = datetime.utcnow().isoformat()
    sermon_dir = UPLOAD_DIR / sermon_id
    sermon_dir.mkdir(parents=True, exist_ok=True)
    destination = sermon_dir / file.filename
    with destination.open("wb") as out_file:
        shutil.copyfileobj(file.file, out_file)

    db.execute(
        """
        INSERT INTO sermons (
            id, sermon_name, series_name, week_or_date, pastor_name,
            status, file_path, original_filename, created_at
        )
        VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)
        """,
        (
            sermon_id,
            sermonName,
            seriesName,
            weekOrDate,
            pastorName,
            "uploaded",
            f"{sermon_id}/{file.filename}",
            file.filename,
            created_at,
        ),
    )
    db.commit()
    init_sermon_state(sermon_id)

    return Sermon(
        id=sermon_id,
        sermonName=sermonName,
        seriesName=seriesName,
        weekOrDate=weekOrDate,
        pastorName=pastorName,
        status="uploaded",
        filePath=f"uploads/{sermon_id}/{file.filename}",
        originalFilename=file.filename,
        createdAt=datetime.fromisoformat(created_at),
    )


@app.get("/sermons", response_model=List[Sermon])
def list_sermons(db=Depends(get_db)) -> List[Sermon]:
    rows = db.execute(
        """
        SELECT
            id, sermon_name, series_name, week_or_date, pastor_name,
            status, file_path, original_filename, created_at
        FROM sermons
        ORDER BY created_at DESC
        """
    ).fetchall()
    return [_row_to_sermon(row) for row in rows]


@app.get("/sermons/{sermon_id}/slides", response_model=List[SlideContent])
def list_sermon_slides(sermon_id: str, db=Depends(get_db)) -> List[SlideContent]:
    presentation = _get_presentation(db, sermon_id)
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        slide_id = f"{sermon_id}:{index}"
        slides.append(
            SlideContent(
                slideId=slide_id,
                slideNumber=index,
                originalText=_extract_slide_text(slide),
            )
        )
    return slides


@app.post(
    "/sermons/{sermon_id}/slides/{slide_number}/analyze",
    response_model=SlideAnalysis,
)
def analyze_slide(
    sermon_id: str, slide_number: int, db=Depends(get_db)
) -> SlideAnalysis:
    if slide_number < 1:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid slide number"
        )

    presentation = _get_presentation(db, sermon_id)
    try:
        slide = presentation.slides[slide_number - 1]
    except IndexError:
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Not found")

    slide_id = f"{sermon_id}:{slide_number}"
    original_text = _extract_slide_text(slide)
    try:
        suggestions = analyze_slide_text(slide_id, original_text)
    except (BedrockAgentError, ValueError, KeyError) as exc:
        raise HTTPException(
            status_code=status.HTTP_502_BAD_GATEWAY,
            detail=f"Bedrock analysis failed: {exc}",
        ) from exc
    analysis = SlideAnalysis(
        slideId=slide_id,
        slideNumber=slide_number,
        originalText=original_text,
        suggestions=suggestions,
    )

    init_sermon_state(sermon_id)
    doc = load_analysis(sermon_id)
    for idx, existing in enumerate(doc.slides):
        if existing.slideId == slide_id:
            doc.slides[idx] = analysis
            break
    else:
        doc.slides.append(analysis)
    save_analysis(doc)

    return analysis


@app.get("/sermons/{sermon_id}/analysis", response_model=AnalysisDocument)
def get_sermon_analysis(sermon_id: str, db=Depends(get_db)) -> AnalysisDocument:
    _ensure_sermon_exists(db, sermon_id)
    init_sermon_state(sermon_id)
    return load_analysis(sermon_id)


@app.get(
    "/sermons/{sermon_id}/slides/{slide_number}/analysis",
    response_model=SlideAnalysis,
)
def get_slide_analysis(
    sermon_id: str, slide_number: int, db=Depends(get_db)
) -> SlideAnalysis:
    if slide_number < 1:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid slide number"
        )

    _ensure_sermon_exists(db, sermon_id)
    init_sermon_state(sermon_id)

    slide_id = f"{sermon_id}:{slide_number}"
    doc = load_analysis(sermon_id)
    for slide in doc.slides:
        if slide.slideId == slide_id:
            return slide

    raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Not found")


@app.post(
    "/sermons/{sermon_id}/slides/{slide_number}/decisions",
    response_model=SlideDecision,
)
def save_slide_decisions(
    sermon_id: str,
    slide_number: int,
    payload: SlideDecisionPayload,
    db=Depends(get_db),
) -> SlideDecision:
    if slide_number < 1:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST, detail="Invalid slide number"
        )

    _ensure_sermon_exists(db, sermon_id)
    init_sermon_state(sermon_id)

    slide_id = f"{sermon_id}:{slide_number}"
    decision = SlideDecision(
        slideId=slide_id,
        slideNumber=slide_number,
        decisions=payload.decisions,
    )

    doc = load_decisions(sermon_id)
    doc.updatedAt = datetime.utcnow()
    for idx, existing in enumerate(doc.slides):
        if existing.slideId == slide_id:
            doc.slides[idx] = decision
            break
    else:
        doc.slides.append(decision)
    save_decisions(doc)

    return decision


@app.get("/sermons/{sermon_id}/decisions", response_model=DecisionsDocument)
def get_sermon_decisions(sermon_id: str, db=Depends(get_db)) -> DecisionsDocument:
    _ensure_sermon_exists(db, sermon_id)
    init_sermon_state(sermon_id)
    return load_decisions(sermon_id)


def _output_pptx_path(sermon_id: str) -> Path:
    return STORAGE_DIR / "sermons" / sermon_id / "output.pptx"


def _replace_in_text_frame(text_frame, replacements: List[tuple[str, str]]) -> None:
    for paragraph in text_frame.paragraphs:
        for original, replacement in replacements:
            if not original or original not in paragraph.text:
                continue
            replaced_in_runs = False
            for run in paragraph.runs:
                if original in run.text:
                    run.text = run.text.replace(original, replacement)
                    replaced_in_runs = True
            if not replaced_in_runs:
                updated = paragraph.text.replace(original, replacement)
                if updated != paragraph.text:
                    paragraph.text = updated


def _apply_text_replacements(slide, replacements: List[tuple[str, str]]) -> None:
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        _replace_in_text_frame(shape.text_frame, replacements)

    try:
        notes_frame = slide.notes_slide.notes_text_frame
    except Exception:
        notes_frame = None

    if notes_frame is not None:
        _replace_in_text_frame(notes_frame, replacements)


@app.post("/sermons/{sermon_id}/generate-updated-pptx")
def generate_updated_pptx(sermon_id: str, db=Depends(get_db)) -> dict:
    _ensure_sermon_exists(db, sermon_id)
    init_sermon_state(sermon_id)

    analysis = load_analysis(sermon_id)
    decisions = load_decisions(sermon_id)
    suggestion_map = {}
    for slide in analysis.slides:
        for suggestion in slide.suggestions:
            suggestion_map[suggestion.id] = suggestion

    presentation = _get_presentation(db, sermon_id)
    for index, slide in enumerate(presentation.slides, start=1):
        slide_id = f"{sermon_id}:{index}"
        decision_block = next(
            (entry for entry in decisions.slides if entry.slideId == slide_id), None
        )
        if not decision_block:
            continue

        replacements = []
        for decision in decision_block.decisions:
            suggestion = suggestion_map.get(decision.suggestionId)
            if not suggestion:
                continue
            if decision.decision == "rejected":
                continue
            if decision.decision == "accepted":
                replacement = suggestion.proposed
            else:
                replacement = decision.finalText or suggestion.proposed
            replacements.append((suggestion.original, replacement))

        if replacements:
            _apply_text_replacements(slide, replacements)

    output_path = _output_pptx_path(sermon_id)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)
    return {"status": "ready"}


@app.get("/sermons/{sermon_id}/download-updated-pptx")
def download_updated_pptx(sermon_id: str, db=Depends(get_db)) -> FileResponse:
    _ensure_sermon_exists(db, sermon_id)
    output_path = _output_pptx_path(sermon_id)
    if not output_path.exists():
        raise HTTPException(status_code=status.HTTP_404_NOT_FOUND, detail="Not found")
    return FileResponse(
        output_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=f"{sermon_id}-updated.pptx",
    )
